from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    
    # Brand Colors
    ELECTRIC_INDIGO = RGBColor(102, 0, 255)
    INK_BLACK = RGBColor(17, 24, 39)
    
    slides_data = [
        {
            "title": "Executive Summary: From Fragmented Services to an Optimized Ecosystem",
            "content": [
                "The Objective: Audit the physical, digital, and operational workflows of Gigabyte Innovation Hub to identify bottlenecks and unlock hidden revenue.",
                "The Current Reality: Despite high technical competence (88% repair conversion) and strong foot traffic (60+ retail leads/month), decentralized systems are causing massive revenue leakage.",
                "The Core Bottlenecks: Decentralized payments, a vacant marketing role, unsecured physical spaces, and siloed departments are currently capping our growth.",
                "The 90-Day Strategy: Centralize front-of-house intake and finance, optimize underutilized physical spaces, and launch standard follow-up systems.",
                "The Ultimate Goal: Double Retail Sales (from 30 to 60/month) and Repair Volume (from 40 to 80/month) without increasing ad spend, purely through operational efficiency."
            ],
            "notes": "Welcome, team. We have spent the last few weeks conducting a deep-dive operational audit... (use the speaker notes drafted in our session)"
        },
        {
            "title": "Audit Methodology: A 360-Degree Ecosystem Diagnosis",
            "content": [
                "The Value Chain Framework: Analyzed Retail, Repair, Academy, and Software across Input, Infrastructure, Attraction, and Output phases.",
                "Physical and Spatial Inspection: Conducted a visual walkthrough of the customer journey, from exterior constraints to internal workstation ergonomics.",
                "Digital and Funnel Audit: Analyzed digital storefronts (Facebook, WhatsApp) to map how online traffic converts to in-store footfall.",
                "Staff Intelligence: Evaluated sales data (lead-to-conversion) and analyzed AI-transcribed meeting logs to uncover the 'ground truth'."
            ],
            "notes": "To figure out exactly how to double our revenue, we didn't rely on guesswork..."
        },
        {
            "title": "Key Findings: The Anatomy of Our Operational Bottlenecks",
            "content": [
                "The Front-of-House & Marketing Void: Decentralized payments and a vacant marketing role causing 50% retail lead leakage.",
                "Supply Chain Vulnerabilities: 52.2% reliance on a single vendor (Supplier B) and localized parts sourcing delays.",
                "Physical & Digital Infrastructure Gaps: High-value inventory is unsecured; software room is currently used as a hardware graveyard.",
                "Product Strategy & Silos: Software build relies on 'Founder's Gut' lacking a roadmap; departments fail to cross-sell (e.g., Retail to Academy)."
            ],
            "notes": "Here is what the data showed. Our staff is working hard, but the system is failing them..."
        },
        {
            "title": "Risk Assessment: The Cost of Operational Inaction",
            "content": [
                "Financial/Accounting Risk: Decentralized payments create massive revenue integrity blind spots.",
                "Asset/Liability Risk: High-ticket retail stock is unsecured (shrinkage risk); Repair Lab lacks ESD protection (liability risk).",
                "Brand/Market Share Risk: Public WhatsApp troubleshooting kills buying momentum; dormant Facebook page cedes market share.",
                "Opportunity Cost: Wasting premium software room real estate as hardware storage."
            ],
            "notes": "Now that we know the bottlenecks, we have to look at what happens if we do nothing..."
        },
        {
            "title": "Target Architecture: Centralized & Optimized Flow",
            "content": [
                "Insert Functional Architectural Map Here",
                "Move from Fragmented/Decentralized to Centralized/Synergistic.",
                "Implement 'The Face' for lead capture and 'The Engine' for financial control."
            ],
            "notes": "This is our new operating system. This map shows how we eliminate bottlenecks."
        },
        {
            "title": "Strategic Recommendations: The Optimization Blueprint",
            "content": [
                "1. Centralize Front-of-House: 'The Face' captures leads & executes follow-ups; 'The Engine' centralizes payments.",
                "2. Asset Monetization: Launch the 'Gigabyte Tech Workspace' (Co-working) to generate daily cash.",
                "3. Supply Chain: Execute sourcing trip to break reliance on Supplier B.",
                "4. Digital Ecosystem: Unified CRM for cross-sells; reactivate Facebook 3x/week.",
                "5. Software Strategy: Require 6-month roadmap and user research."
            ],
            "notes": "So, how do we fix this? Centralize the desk, monetize the storage room, and build a CRM."
        },
        {
            "title": "The 90-Day Execution Roadmap",
            "content": [
                "Phase 1 (Days 1-30): Stop the Bleeding (Split roles, No-Drop Follow-up SOP, Clean software room).",
                "Phase 2 (Days 31-60): Secure & Optimize (Co-working launch, Glass cabinets, Sourcing trip).",
                "Phase 3 (Days 61-90): Scale & Connect (Unified CRM, Automated cross-selling, Software roadmaps)."
            ],
            "notes": "Ideas don't generate revenue; execution does. We are doing this systematically."
        },
        {
            "title": "Questions and Open Discussion",
            "content": [
                "Defending the Strategy:",
                "- Why split the front desk? (Revenue leakage)",
                "- Why ESD mats/Bin & Tag? (Liability protection)",
                "- Why Co-working? (Asset utilization)",
                "- Why roadmap? (Scaling)"
            ],
            "notes": "Thank you. I am requesting approval to initiate Phase 1 by June 2026."
        }
    ]

    for slide_data in slides_data:
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title = slide.shapes.title
        title.text = slide_data['title']
        
        # Content
        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        
        for item in slide_data['content']:
            p = tf.add_paragraph()
            p.text = item
            p.level = 0
            
        # Add Notes
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = slide_data['notes']

    prs.save('Gigabyte_Operational_Audit_Presentation.pptx')
    print("Presentation generated successfully!")

if __name__ == "__main__":
    create_presentation()
